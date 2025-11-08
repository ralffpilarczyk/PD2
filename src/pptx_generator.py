"""
PowerPoint generator for OnePageProfile (OPP.py)
Converts markdown profiles to formatted PowerPoint slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
from pathlib import Path
from typing import Dict, List
from datetime import datetime


def _get_opp_version():
    """Get version from OPP.py by reading the file directly (avoid circular import)"""
    opp_path = Path(__file__).parent.parent / "OPP.py"
    with open(opp_path, 'r') as f:
        for line in f:
            if line.startswith('__opp_version__'):
                # Extract version string from line like: __opp_version__ = "1.0"
                return line.split('=')[1].strip().strip('"').strip("'")
    return "1.0"  # Fallback


OPP_VERSION = _get_opp_version()


# Color scheme (matching PD2)
DARK_BLUE = RGBColor(45, 90, 135)  # #2d5a87 - titles and bold keywords
DARK_GREY = RGBColor(74, 85, 104)  # #4a5568 - regular text


def parse_markdown_profile(md_path: str) -> Dict:
    """Parse markdown profile into structured data

    Returns:
        {
            'title': str,
            'subtitle': str,
            'sections': {
                'Company Overview': [bullet1, bullet2, ...],
                'Competitive Positioning': [...],
                'Financial KPIs': [...],
                'Strategic Considerations': [...]
            }
        }
    """
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    result = {
        'title': '',
        'subtitle': '',
        'sections': {}
    }

    # Extract title (# Title)
    title_match = re.search(r'^#\s+(.+?)$', content, re.MULTILINE)
    if title_match:
        result['title'] = title_match.group(1).strip()

    # Extract subtitle (first line after title, before first ##)
    subtitle_match = re.search(r'^#\s+.+?\n(.+?)(?=\n##)', content, re.MULTILINE | re.DOTALL)
    if subtitle_match:
        result['subtitle'] = subtitle_match.group(1).strip()

    # Extract sections and their bullets
    section_pattern = r'##\s+(.+?)\n(.*?)(?=\n##|\Z)'
    for match in re.finditer(section_pattern, content, re.MULTILINE | re.DOTALL):
        section_name = match.group(1).strip()
        section_content = match.group(2).strip()

        # Extract bullet points (lines starting with *)
        bullets = []
        for line in section_content.split('\n'):
            line = line.strip()
            if line.startswith('*'):
                # Remove leading * and whitespace
                bullet = line.lstrip('*').strip()
                if bullet:  # Only add non-empty bullets
                    bullets.append(bullet)

        if bullets:  # Only add sections with content
            result['sections'][section_name] = bullets

    return result


def create_profile_pptx(md_path: str, company_name: str, timestamp: str, version_suffix: str = "", profile_type: str = "default") -> str:
    """Create PowerPoint presentation from markdown profile

    Args:
        md_path: Path to step4_final.md
        company_name: Company name for filename
        timestamp: Timestamp for filename
        version_suffix: Optional version suffix (e.g., "_v1", "_v2", "_v3")
        profile_type: Profile type ("default" or "custom") for filename prefix

    Returns:
        Path to generated PPTX file
    """
    # Parse markdown
    data = parse_markdown_profile(md_path)

    # Create presentation
    prs = Presentation()

    # Set slide size to A4 landscape (11.69" x 8.27")
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title and subtitle at top
    _add_title_area(slide, data['title'], data['subtitle'])

    # Add 2x2 grid of section boxes
    _add_2x2_boxes(slide, data['sections'])

    # Add footnote at bottom
    # _add_footnote(slide)  # Removed per user request

    # Save to ProfileFiles
    output_dir = Path("ProfileFiles")
    output_dir.mkdir(exist_ok=True)

    # Clean company name for filename
    clean_name = "".join(c for c in company_name if c.isalnum() or c in (' ', '_', '-')).strip()
    clean_name = clean_name.replace(' ', '_')

    # Add prefix for custom profiles
    prefix = "Custom_" if profile_type == "custom" else ""
    output_path = output_dir / f"{prefix}{clean_name}_{timestamp}{version_suffix}.pptx"
    prs.save(str(output_path))

    return str(output_path)


def _add_title_area(slide, title: str, subtitle: str):
    """Add title and subtitle at top of slide"""
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),   # left
        Inches(0.5),   # top
        Inches(10.69), # width (11.69 - 0.5 - 0.5)
        Inches(0.7)    # height
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.name = 'Arial'
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5),   # left
        Inches(1.053), # top (moved up 5mm = 0.197")
        Inches(10.69), # width
        Inches(0.6)    # height
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.name = 'Arial'
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.bold = True
    subtitle_frame.paragraphs[0].font.color.rgb = DARK_BLUE


def _add_2x2_boxes(slide, sections: Dict[str, List[str]]):
    """Add 2x2 grid of section boxes

    Layout:
        Top-left: Company Overview
        Top-right: Competitive Positioning
        Bottom-left: Financial KPIs
        Bottom-right: Strategic Considerations
    """
    # Box dimensions
    box_width = Inches(5.195)   # (10.69 - 0.3) / 2
    box_height = Inches(2.635)  # (5.57 - 0.3) / 2

    # Positions: (section_name, left, top)
    # All boxes moved up by 10mm total (0.394")
    positions = [
        ('Company Overview', Inches(0.5), Inches(1.806)),
        ('Competitive Positioning', Inches(6.195), Inches(1.806)),
        ('Financial KPIs', Inches(0.5), Inches(4.741)),
        ('Strategic Considerations', Inches(6.195), Inches(4.741))
    ]

    for section_name, left, top in positions:
        if section_name in sections:
            _add_section_box(slide, section_name, sections[section_name], left, top, box_width, box_height)


def _add_section_box(slide, section_name: str, bullets: List[str], left, top, width, height):
    """Add a single section box with title and bullets"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)

    # Section title
    p = text_frame.paragraphs[0]
    p.text = section_name
    p.font.name = 'Arial'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(6)

    # Bullets
    for bullet_text in bullets:
        p = text_frame.add_paragraph()
        # Enable native PowerPoint bullets (handles hanging indent automatically)
        _enable_bullet(p)
        # Add formatted content (no manual "• " needed - PowerPoint adds it)
        _format_bullet_text(p, bullet_text)
        p.space_after = Pt(3)


def _format_bullet_text(paragraph, text: str):
    """Format bullet text with bold keywords

    Parse markdown **bold** syntax and apply formatting:
    - Regular text: Arial 10pt, dark grey
    - Bold text: Arial Bold 10pt, dark blue
    """
    # Fix malformed bold: "Label**: " → "**Label**: "
    text = re.sub(r'([^*:]+?)\*\*:', r'**\1**:', text)

    # Split text by **bold** markers
    parts = re.split(r'(\*\*.*?\*\*)', text)

    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            # Bold keyword
            keyword = part[2:-2]  # Remove ** markers
            run = paragraph.add_run()
            run.text = keyword
            run.font.name = 'Arial'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = DARK_BLUE
        elif part:  # Regular text
            run = paragraph.add_run()
            run.text = part
            run.font.name = 'Arial'
            run.font.size = Pt(10)
            run.font.color.rgb = DARK_GREY


def _enable_bullet(paragraph):
    """Enable PowerPoint's native bullet formatting on a paragraph

    This sets up standard bullet formatting with proper spacing:
    - Bullet character: •
    - Space between bullet and text
    - Wrapped lines align with first line text (not the bullet)
    """
    from lxml import etree

    # Access the paragraph properties element
    pPr = paragraph._element.get_or_add_pPr()

    # Set indentation (in EMUs: 1 inch = 914400 EMUs)
    # marL: where text starts (0.125 inches = 114300 EMUs)
    # indent: where bullet is, relative to marL (negative for hanging)
    pPr.set('marL', '114300')    # Left margin: text starts here
    pPr.set('indent', '-114300')  # Hanging indent: bullet positioned 0.125" left of text

    # Create bullet character element
    buChar = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    if buChar is None:
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    buChar.set('char', '•')

    # Set bullet font
    buFont = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
    if buFont is None:
        buFont = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
    buFont.set('typeface', 'Arial')


def _add_footnote(slide):
    """Add footnote at bottom of slide"""
    # Format current date as DD-MMM-YY
    date_str = datetime.now().strftime("%d-%b-%y")

    footnote_box = slide.shapes.add_textbox(
        Inches(0.5),   # left
        Inches(8.05),  # bottom of page (slide height is 8.27")
        Inches(10.69), # width
        Inches(0.2)    # height
    )
    footnote_frame = footnote_box.text_frame
    footnote_frame.text = f"Note: Generated with OnePageProfile v{OPP_VERSION} as of {date_str}."
    footnote_frame.paragraphs[0].font.name = 'Arial'
    footnote_frame.paragraphs[0].font.size = Pt(7)
    footnote_frame.paragraphs[0].font.color.rgb = DARK_BLUE
